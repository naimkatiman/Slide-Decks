"""
Compile 10 HTML slide files into a single PowerPoint presentation.
Screenshots each HTML at 1280x720 and embeds as full-slide images.
Adds PowerPoint slide transitions for animation effect.
"""
import os
import sys
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu

SLIDE_DIR = Path(r"d:\Personal\Solution Engineer Coin Gecko\Slide Decks")
OUTPUT_FILE = SLIDE_DIR / "CoinGecko_DogePay_Presentation.pptx"
SCREENSHOT_DIR = SLIDE_DIR / "_screenshots"
SLIDE_WIDTH = 1280
SLIDE_HEIGHT = 720
NUM_SLIDES = 10


def screenshot_slides():
    """Take screenshots of all HTML slides using Playwright."""
    SCREENSHOT_DIR.mkdir(exist_ok=True)

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": SLIDE_WIDTH, "height": SLIDE_HEIGHT})

        for i in range(1, NUM_SLIDES + 1):
            html_file = SLIDE_DIR / f"{i}.html"
            if not html_file.exists():
                print(f"WARNING: {html_file} not found, skipping")
                continue

            file_url = html_file.as_uri()
            page.goto(file_url)
            # Wait for animations to settle
            page.wait_for_timeout(2000)

            screenshot_path = SCREENSHOT_DIR / f"slide_{i:02d}.png"
            page.screenshot(path=str(screenshot_path), full_page=False)
            print(f"Captured slide {i}: {screenshot_path}")

        browser.close()


def compile_pptx():
    """Compile screenshots into a PowerPoint presentation with transitions."""
    prs = Presentation()

    # Set slide dimensions to 16:9 (1280x720 -> 13.333 x 7.5 inches at 96 DPI)
    prs.slide_width = Emu(12192000)   # 13.333 inches
    prs.slide_height = Emu(6858000)   # 7.5 inches

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for i in range(1, NUM_SLIDES + 1):
        screenshot_path = SCREENSHOT_DIR / f"slide_{i:02d}.png"
        if not screenshot_path.exists():
            print(f"WARNING: {screenshot_path} not found, skipping")
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Add image to fill the entire slide
        slide.shapes.add_picture(
            str(screenshot_path),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # Add slide transition (fade effect)
        from pptx.oxml.ns import qn

        transition = slide._element.makeelement(
            qn('p:transition'), {'spd': 'med', 'advClick': '1'}
        )
        fade = transition.makeelement(qn('p:fade'), {})
        transition.append(fade)
        # Insert transition after cSld element
        cSld = slide._element.find(qn('p:cSld'))
        cSld.addnext(transition)

        print(f"Added slide {i} to presentation")

    prs.save(str(OUTPUT_FILE))
    print(f"\nPresentation saved: {OUTPUT_FILE}")
    print(f"Total slides: {NUM_SLIDES}")


def main():
    print("=" * 50)
    print("CoinGecko x DogePay - PPTX Compiler")
    print("=" * 50)

    print("\n[1/2] Capturing slide screenshots...")
    screenshot_slides()

    print("\n[2/2] Compiling PowerPoint presentation...")
    compile_pptx()

    print("\nDone!")


if __name__ == "__main__":
    main()
